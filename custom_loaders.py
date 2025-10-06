from langchain.schema import Document
from docx import Document as DocxDocument
from pptx import Presentation

class DocxLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        doc = DocxDocument(self.file_path)
        full_text = []
        for paragraph in doc.paragraphs:
            full_text.append(paragraph.text)
        return [Document(page_content="\n".join(full_text))]

class PptxLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        prs = Presentation(self.file_path)
        slides_content = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            slides_content.append("\n".join(slide_text))
        return [Document(page_content="\n".join(slides_content))]
